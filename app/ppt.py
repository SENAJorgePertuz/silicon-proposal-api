# backend/app/ppt.py
"""
Módulo encargado de la manipulación de la plantilla PPTX.

Contiene funciones para:
- Reemplazar placeholders en el texto.
- Leer tags de control desde las notas de las diapositivas.
- Generar una presentación personalizada en memoria.

Todo el procesamiento se hace en memoria (BytesIO); no se escribe en disco.
"""
from pptx import Presentation
from io import BytesIO
from typing import Dict, Any
import re

def replace_placeholders_in_text(text_frame, replacements: Dict[str, str]) -> None:
    """
    Reemplaza todos los placeholders definidos en `replacements` dentro de un `text_frame`.
    
    Recorre cada párrafo y cada "run" de texto (fragmento con formato uniforme)
    para evitar romper el estilo al reemplazar.
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                for key, value in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, value)

def get_slide_tags(notes_text: str) -> list:
    """
    Extrae los tags de control desde el texto de las notas de una diapositiva.
    
    Formato esperado en las notas: [[tag:nombre_del_tag]]
    Ejemplo: [[tag:about_eic_accelerator]]
    
    Devuelve una lista de nombres de tags (sin el prefijo "tag:").
    """
    pattern = r"\[\[tag:([a-zA-Z0-9_]+)\]\]"
    return re.findall(pattern, notes_text)

def generate_presentation(
    template_path: str,
    replacements: Dict[str, str],
    slide_toggles: Dict[str, bool]
) -> BytesIO:
    """
    Genera una presentación PPTX personalizada a partir de una plantilla.
    
    Pasos:
    1. Carga la plantilla desde `template_path`.
    2. Identifica diapositivas a eliminar según los tags en sus notas y `slide_toggles`.
    3. Elimina las diapositivas marcadas (de atrás hacia adelante para preservar índices).
    4. Reemplaza todos los placeholders en el texto visible y en las notas.
    5. Devuelve el resultado como un buffer en memoria (BytesIO).
    
    Nota: Si un tag no está en `slide_toggles`, la diapositiva se incluye por defecto.
    """
    prs = Presentation(template_path)
    slides_to_remove = []

    # Paso 1: Identificar diapositivas a eliminar
    for i, slide in enumerate(prs.slides):
        
        # Leer el contenido de las notas (si existen)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text

        # Extraer tags de las notas
        tags = get_slide_tags(notes)
        should_remove = False

        # Si algún tag está desactivado en slide_toggles, marcar para eliminar
        for tag in tags:
            if not slide_toggles.get(tag, True):  # por defecto: incluir si no está en toggles
                should_remove = True
                break

        if should_remove:
            slides_to_remove.append(i)

    # Paso 2: Eliminar diapositivas (de atrás hacia adelante para no alterar índices)
    for i in reversed(slides_to_remove):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # Paso 3: Reemplazar placeholders en todas las diapositivas restantes
    for slide in prs.slides:
        # En el cuerpo de la diapositiva
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            replace_placeholders_in_text(shape.text_frame, replacements)

        # También en las notas (útil para depuración o futuras versiones)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            replace_placeholders_in_text(slide.notes_slide.notes_text_frame, replacements)

    # Paso 4: Guardar en memoria y devolver
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output